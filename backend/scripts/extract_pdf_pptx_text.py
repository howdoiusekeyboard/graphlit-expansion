"""Extract text from PDF and PPTX assets into Markdown for review.

Usage:
    python scripts/extract_pdf_pptx_text.py --out new/premidsem_assets_extracted.md
"""

from __future__ import annotations

import argparse
from pathlib import Path

from pypdf import PdfReader
from pptx import Presentation


def extract_pdf_text(path: Path) -> list[str]:
    reader = PdfReader(str(path))
    lines: list[str] = []
    lines.append("---")
    lines.append(f"## {path.name}")
    lines.append("")
    lines.append(f"**Type:** PDF  |  **Pages:** {len(reader.pages)}")
    lines.append("")
    for i, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if not text:
            continue
        lines.append(f"### Page {i}")
        lines.append("")
        # Keep it readable; preserve line breaks but collapse excessive whitespace
        for raw_line in text.splitlines():
            t = raw_line.strip()
            if t:
                lines.append(t)
        lines.append("")
    return lines


def extract_pptx_text(path: Path) -> list[str]:
    prs = Presentation(str(path))
    lines: list[str] = []
    lines.append("---")
    lines.append(f"## {path.name}")
    lines.append("")
    lines.append(f"**Type:** PPTX  |  **Slides:** {len(prs.slides)}")
    lines.append("")

    for s_i, slide in enumerate(prs.slides, start=1):
        slide_lines: list[str] = []
        for shape in slide.shapes:
            if not hasattr(shape, "has_text_frame"):
                continue
            if not shape.has_text_frame:
                continue
            text = (shape.text or "").strip()
            if text:
                slide_lines.append(text)

        if not slide_lines:
            continue

        lines.append(f"### Slide {s_i}")
        lines.append("")
        for t in slide_lines:
            # Separate blocks for readability
            lines.append(t)
            lines.append("")

    return lines


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--out", type=str, default="new/premidsem_assets_extracted.md")
    args = parser.parse_args()

    out_path = Path(args.out)
    new_dir = Path("new")

    pdfs = sorted(new_dir.glob("*.pdf"))
    pptxs = sorted(new_dir.glob("*.pptx"))

    lines: list[str] = []
    lines.append("# Premidsem PDF/PPTX extraction")
    lines.append("")
    lines.append("Source: `new/`")
    lines.append("")

    if not pdfs and not pptxs:
        lines.append("_No PDF/PPTX files found in `new/`._")
        out_path.write_text("\n".join(lines), encoding="utf-8")
        return 0

    for p in pdfs:
        lines.extend(extract_pdf_text(p))

    for p in pptxs:
        lines.extend(extract_pptx_text(p))

    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_path.write_text("\n".join(lines), encoding="utf-8")
    print(f"Wrote extraction to {out_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())



